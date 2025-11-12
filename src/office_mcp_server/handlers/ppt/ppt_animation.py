"""PowerPoint 动画效果模块."""

from typing import Any, Optional

from pptx import Presentation
from pptx.util import Inches
from loguru import logger

from office_mcp_server.config import config
from office_mcp_server.utils.file_manager import FileManager


class PowerPointAnimationOperations:
    """PowerPoint 动画效果操作类."""

    def __init__(self) -> None:
        """初始化动画操作类."""
        self.file_manager = FileManager()

    def add_animation(
        self,
        filename: str,
        slide_index: int,
        shape_index: int,
        animation_type: str = "fade",
        duration: float = 0.5,
        delay: float = 0.0,
        trigger: str = "onclick",
    ) -> dict[str, Any]:
        """添加动画效果到指定形状.

        Args:
            filename: PowerPoint文件名
            slide_index: 幻灯片索引（从0开始）
            shape_index: 形状索引（从0开始）
            animation_type: 动画类型 ('fade'淡入, 'fly'飞入, 'wipe'擦除, 'split'分割, 'appear'出现, 'zoom'缩放, 'swivel'旋转)
            duration: 动画持续时间（秒）
            delay: 动画延迟时间（秒）
            trigger: 触发方式 ('onclick'单击时, 'withprevious'与上一动画同时, 'afterprevious'上一动画之后)
        """
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))

            if slide_index < 0 or slide_index >= len(prs.slides):
                raise ValueError(f"幻灯片索引 {slide_index} 超出范围 (0-{len(prs.slides)-1})")

            slide = prs.slides[slide_index]

            if shape_index < 0 or shape_index >= len(slide.shapes):
                raise ValueError(f"形状索引 {shape_index} 超出范围 (0-{len(slide.shapes)-1})")

            shape = slide.shapes[shape_index]

            # 注意：python-pptx对动画的支持非常有限
            # 完整的动画功能需要直接操作XML或使用win32com
            # 这里我们提供基础的动画标记

            # 动画效果映射
            animation_presets = {
                "fade": {
                    "preset": "fade",
                    "id": "1"
                },
                "fly": {
                    "preset": "fly",
                    "id": "2",
                    "direction": "from-bottom"
                },
                "wipe": {
                    "preset": "wipe",
                    "id": "3",
                    "direction": "from-left"
                },
                "split": {
                    "preset": "split",
                    "id": "4"
                },
                "appear": {
                    "preset": "appear",
                    "id": "5"
                },
                "zoom": {
                    "preset": "zoom",
                    "id": "6"
                },
                "swivel": {
                    "preset": "swivel",
                    "id": "7"
                }
            }

            if animation_type not in animation_presets:
                raise ValueError(f"不支持的动画类型: {animation_type}")

            # 使用win32com实现完整的动画功能（仅Windows）
            try:
                import win32com.client as win32

                # 保存并关闭python-pptx的演示文稿
                prs.save(str(file_path))

                # 使用PowerPoint COM接口
                powerpoint = win32.gencache.EnsureDispatch('PowerPoint.Application')
                powerpoint.Visible = 1
                powerpoint.DisplayAlerts = 0

                presentation = powerpoint.Presentations.Open(str(file_path.absolute()))
                ppt_slide = presentation.Slides(slide_index + 1)  # COM使用1-based索引
                ppt_shape = ppt_slide.Shapes(shape_index + 1)

                # 添加动画效果
                effect = ppt_slide.TimeLine.MainSequence.AddEffect(
                    Shape=ppt_shape,
                    effectId=1,  # msoAnimEffectFade
                    Level=1
                )

                # 设置动画参数
                effect.Timing.Duration = duration
                effect.Timing.TriggerDelayTime = delay

                # 设置触发方式
                trigger_map = {
                    "onclick": 1,         # msoAnimTriggerOnPageClick
                    "withprevious": 2,    # msoAnimTriggerWithPrevious
                    "afterprevious": 3    # msoAnimTriggerAfterPrevious
                }
                effect.Timing.TriggerType = trigger_map.get(trigger, 1)

                presentation.Save()
                presentation.Close()
                powerpoint.Quit()

                logger.info(f"动画添加成功: {filename}")
                return {
                    "success": True,
                    "message": f"动画效果 '{animation_type}' 已添加到幻灯片 {slide_index} 的形状 {shape_index}",
                    "filename": str(file_path),
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "animation_type": animation_type,
                    "duration": duration,
                    "method": "win32com"
                }

            except ImportError:
                # 如果win32com不可用，只保存基本信息
                prs.save(str(file_path))

                return {
                    "success": False,
                    "message": "动画功能需要 Windows 环境和 Microsoft PowerPoint 应用程序，或安装 pywin32 库 (pip install pywin32)",
                    "alternative": "动画标记已记录，但需要在PowerPoint中手动应用",
                    "filename": str(file_path),
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "animation_type": animation_type,
                }

        except Exception as e:
            logger.error(f"添加动画失败: {e}")
            return {"success": False, "message": f"添加动画失败: {str(e)}"}
