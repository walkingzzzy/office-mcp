"""PowerPoint 基础操作模块."""

from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Inches
from loguru import logger

from office_mcp_server.config import config
from office_mcp_server.utils.file_manager import FileManager


class PowerPointBasicOperations:
    """PowerPoint 基础操作类."""

    def __init__(self) -> None:
        """初始化基础操作类."""
        self.file_manager = FileManager()

    def create_presentation(
        self, filename: str, title: str = "", template_path: Optional[str] = None
    ) -> dict[str, Any]:
        """创建 PowerPoint 演示文稿.

        Args:
            filename: 文件名
            title: 演示标题（可选）
            template_path: 模板文件路径（可选，如果提供则基于模板创建）
        """
        try:
            file_path = self.file_manager.validate_file_path(filename)
            self.file_manager.validate_file_extension(filename, [".pptx"])

            output_path = config.paths.output_dir / file_path.name
            self.file_manager.ensure_directory(output_path.parent)

            # 如果提供了模板路径，基于模板创建
            if template_path:
                template_file = Path(template_path)
                if not template_file.exists():
                    # 尝试在 output_dir 中查找模板
                    template_file = config.paths.output_dir / template_path
                    if not template_file.exists():
                        raise FileNotFoundError(f"模板文件不存在: {template_path}")

                prs = Presentation(str(template_file))
                logger.info(f"基于模板创建演示文稿: {template_file}")
            else:
                prs = Presentation()

            # 添加标题页（仅在非模板模式下）
            if title and not template_path:
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title_shape = slide.shapes.title
                title_shape.text = title

            prs.save(str(output_path))

            logger.info(f"PowerPoint 演示文稿创建成功: {output_path}")
            return {
                "success": True,
                "message": "PowerPoint 演示文稿创建成功",
                "filename": str(output_path),
                "title": title,
                "from_template": bool(template_path),
            }

        except Exception as e:
            logger.error(f"创建 PowerPoint 演示文稿失败: {e}")
            return {"success": False, "message": f"创建失败: {str(e)}"}

    def add_slide(
        self, filename: str, layout_index: int = 1, title: str = ""
    ) -> dict[str, Any]:
        """添加幻灯片."""
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))
            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)

            # 设置标题
            if title and slide.shapes.title:
                slide.shapes.title.text = title

            prs.save(str(file_path))

            logger.info(f"幻灯片添加成功: {file_path}")
            return {
                "success": True,
                "message": "幻灯片添加成功",
                "filename": str(file_path),
                "slide_count": len(prs.slides),
            }

        except Exception as e:
            logger.error(f"添加幻灯片失败: {e}")
            return {"success": False, "message": f"添加失败: {str(e)}"}

    def delete_slide(self, filename: str, slide_index: int) -> dict[str, Any]:
        """删除幻灯片."""
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))

            if slide_index >= len(prs.slides):
                raise ValueError(f"幻灯片索引 {slide_index} 超出范围")

            # 删除幻灯片
            rId = prs.slides._sldIdLst[slide_index].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[slide_index]

            prs.save(str(file_path))

            logger.info(f"幻灯片删除成功: {file_path}")
            return {
                "success": True,
                "message": "幻灯片删除成功",
                "filename": str(file_path),
                "remaining_slides": len(prs.slides),
            }

        except Exception as e:
            logger.error(f"删除幻灯片失败: {e}")
            return {"success": False, "message": f"删除失败: {str(e)}"}

    def move_slide(
        self, filename: str, from_index: int, to_index: int
    ) -> dict[str, Any]:
        """移动幻灯片位置."""
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))

            slide_count = len(prs.slides)
            if from_index >= slide_count or to_index >= slide_count:
                raise ValueError(f"幻灯片索引超出范围")

            # 移动幻灯片
            slides = list(prs.slides._sldIdLst)
            slide = slides.pop(from_index)
            slides.insert(to_index, slide)

            # 更新幻灯片列表
            prs.slides._sldIdLst.clear()
            for s in slides:
                prs.slides._sldIdLst.append(s)

            prs.save(str(file_path))

            logger.info(f"幻灯片移动成功: {file_path}")
            return {
                "success": True,
                "message": f"幻灯片从位置 {from_index} 移动到位置 {to_index}",
                "filename": str(file_path),
            }

        except Exception as e:
            logger.error(f"移动幻灯片失败: {e}")
            return {"success": False, "message": f"移动失败: {str(e)}"}

    def duplicate_slide(self, filename: str, slide_index: int) -> dict[str, Any]:
        """复制幻灯片."""
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))

            if slide_index >= len(prs.slides):
                raise ValueError(f"幻灯片索引 {slide_index} 超出范围")

            # 获取源幻灯片
            source_slide = prs.slides[slide_index]

            # 创建新幻灯片
            blank_slide_layout = source_slide.slide_layout
            dest_slide = prs.slides.add_slide(blank_slide_layout)

            # 复制所有形状
            for shape in source_slide.shapes:
                el = shape.element
                newel = el.__class__(el)
                dest_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

            prs.save(str(file_path))

            logger.info(f"幻灯片复制成功: {file_path}")
            return {
                "success": True,
                "message": f"幻灯片 {slide_index} 已复制",
                "filename": str(file_path),
                "total_slides": len(prs.slides),
            }

        except Exception as e:
            logger.error(f"复制幻灯片失败: {e}")
            return {"success": False, "message": f"复制失败: {str(e)}"}

    def get_presentation_info(self, filename: str) -> dict[str, Any]:
        """获取演示文稿信息."""
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))
            slide_count = len(prs.slides)

            logger.info(f"获取演示文稿信息成功: {file_path}")
            return {
                "success": True,
                "filename": str(file_path),
                "slide_count": slide_count,
            }

        except Exception as e:
            logger.error(f"获取演示文稿信息失败: {e}")
            return {"success": False, "message": f"获取失败: {str(e)}"}
