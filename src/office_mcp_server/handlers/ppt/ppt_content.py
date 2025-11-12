"""PowerPoint 内容操作模块."""

from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Inches
from loguru import logger

from office_mcp_server.config import config
from office_mcp_server.utils.file_manager import FileManager


class PowerPointContentOperations:
    """PowerPoint 内容操作类."""

    def __init__(self) -> None:
        """初始化内容操作类."""
        self.file_manager = FileManager()

    def add_text(
        self,
        filename: str,
        slide_index: int,
        text: str,
        left_inches: float = 1.0,
        top_inches: float = 1.0,
        width_inches: float = 8.0,
        height_inches: float = 1.0,
    ) -> dict[str, Any]:
        """添加文本框."""
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))

            if slide_index >= len(prs.slides):
                raise ValueError(f"幻灯片索引 {slide_index} 超出范围")

            slide = prs.slides[slide_index]

            # 添加文本框
            left = Inches(left_inches)
            top = Inches(top_inches)
            width = Inches(width_inches)
            height = Inches(height_inches)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = text

            prs.save(str(file_path))

            logger.info(f"文本框添加成功: {file_path}")
            return {
                "success": True,
                "message": "文本框添加成功",
                "filename": str(file_path),
            }

        except Exception as e:
            logger.error(f"添加文本框失败: {e}")
            return {"success": False, "message": f"添加失败: {str(e)}"}

    def add_image(
        self,
        filename: str,
        slide_index: int,
        image_path: str,
        left_inches: float = 1.0,
        top_inches: float = 1.0,
        width_inches: Optional[float] = None,
    ) -> dict[str, Any]:
        """添加图片."""
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            # 验证图片文件
            img_path = Path(image_path)
            if not img_path.exists():
                raise FileNotFoundError(f"图片文件不存在: {image_path}")

            prs = Presentation(str(file_path))

            if slide_index >= len(prs.slides):
                raise ValueError(f"幻灯片索引 {slide_index} 超出范围")

            slide = prs.slides[slide_index]

            # 添加图片
            left = Inches(left_inches)
            top = Inches(top_inches)

            if width_inches:
                width = Inches(width_inches)
                slide.shapes.add_picture(str(img_path), left, top, width=width)
            else:
                slide.shapes.add_picture(str(img_path), left, top)

            prs.save(str(file_path))

            logger.info(f"图片添加成功: {file_path}")
            return {
                "success": True,
                "message": "图片添加成功",
                "filename": str(file_path),
            }

        except Exception as e:
            logger.error(f"添加图片失败: {e}")
            return {"success": False, "message": f"添加失败: {str(e)}"}

    def add_table(
        self,
        filename: str,
        slide_index: int,
        rows: int,
        cols: int,
        data: Optional[list[list[str]]] = None,
    ) -> dict[str, Any]:
        """添加表格."""
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))

            if slide_index >= len(prs.slides):
                raise ValueError(f"幻灯片索引 {slide_index} 超出范围")

            slide = prs.slides[slide_index]

            # 添加表格
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(3.0)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # 填充数据
            if data:
                for i, row_data in enumerate(data):
                    if i >= rows:
                        break
                    for j, cell_data in enumerate(row_data):
                        if j >= cols:
                            break
                        table.cell(i, j).text = str(cell_data)

            prs.save(str(file_path))

            logger.info(f"表格添加成功: {file_path}")
            return {
                "success": True,
                "message": "表格添加成功",
                "filename": str(file_path),
                "rows": rows,
                "cols": cols,
            }

        except Exception as e:
            logger.error(f"添加表格失败: {e}")
            return {"success": False, "message": f"添加失败: {str(e)}"}
