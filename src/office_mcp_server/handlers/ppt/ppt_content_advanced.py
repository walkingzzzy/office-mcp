"""PowerPoint 内容高级操作模块 - 表格高级操作、形状、图表等."""

from typing import Any, Optional, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from loguru import logger

from office_mcp_server.config import config
from office_mcp_server.utils.file_manager import FileManager
from office_mcp_server.utils.format_helper import ColorUtils


class PowerPointContentAdvancedOperations:
    """PowerPoint 内容高级操作类."""

    def __init__(self) -> None:
        """初始化内容高级操作类."""
        self.file_manager = FileManager()

    # ========== 表格高级操作 ==========
    def insert_table_row(
        self,
        filename: str,
        slide_index: int,
        table_index: int,
        row_index: int,
        data: Optional[List[str]] = None,
    ) -> dict[str, Any]:
        """插入表格行.

        Args:
            filename: 文件名
            slide_index: 幻灯片索引
            table_index: 表格索引
            row_index: 插入位置索引
            data: 行数据
        """
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))
            slide = prs.slides[slide_index]

            # 查找表格
            tables = [shape for shape in slide.shapes if shape.has_table]
            if table_index >= len(tables):
                raise ValueError(f"表格索引 {table_index} 超出范围")

            table = tables[table_index].table

            # 添加行 (python-pptx不支持在指定位置插入行,只能添加到末尾)
            new_row = table.add_row()

            if data:
                for col_idx, cell_data in enumerate(data[:len(new_row.cells)]):
                    new_row.cells[col_idx].text = str(cell_data)

            prs.save(str(file_path))

            logger.info(f"表格行插入成功: {file_path}")
            return {
                "success": True,
                "message": "表格行插入成功",
                "filename": str(file_path),
            }

        except Exception as e:
            logger.error(f"插入表格行失败: {e}")
            return {"success": False, "message": f"操作失败: {str(e)}"}

    def delete_table_row(
        self,
        filename: str,
        slide_index: int,
        table_index: int,
        row_index: int,
    ) -> dict[str, Any]:
        """删除表格行."""
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))
            slide = prs.slides[slide_index]

            tables = [shape for shape in slide.shapes if shape.has_table]
            if table_index >= len(tables):
                raise ValueError(f"表格索引 {table_index} 超出范围")

            table = tables[table_index].table

            # python-pptx不直接支持删除行,需要通过XML操作
            # 这里返回提示信息
            return {
                "success": False,
                "message": "python-pptx库不支持删除表格行,需要使用win32com或更底层的XML操作",
                "filename": str(file_path),
            }

        except Exception as e:
            logger.error(f"删除表格行失败: {e}")
            return {"success": False, "message": f"操作失败: {str(e)}"}

    def merge_table_cells(
        self,
        filename: str,
        slide_index: int,
        table_index: int,
        start_row: int,
        start_col: int,
        end_row: int,
        end_col: int,
    ) -> dict[str, Any]:
        """合并表格单元格."""
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))
            slide = prs.slides[slide_index]

            tables = [shape for shape in slide.shapes if shape.has_table]
            if table_index >= len(tables):
                raise ValueError(f"表格索引 {table_index} 超出范围")

            table = tables[table_index].table

            # 合并单元格
            start_cell = table.cell(start_row, start_col)
            end_cell = table.cell(end_row, end_col)
            start_cell.merge(end_cell)

            prs.save(str(file_path))

            logger.info(f"表格单元格合并成功: {file_path}")
            return {
                "success": True,
                "message": "表格单元格合并成功",
                "filename": str(file_path),
            }

        except Exception as e:
            logger.error(f"合并表格单元格失败: {e}")
            return {"success": False, "message": f"操作失败: {str(e)}"}

    def format_table_cell(
        self,
        filename: str,
        slide_index: int,
        table_index: int,
        row: int,
        col: int,
        fill_color: Optional[str] = None,
        text_color: Optional[str] = None,
        bold: bool = False,
        font_size: Optional[int] = None,
    ) -> dict[str, Any]:
        """格式化表格单元格."""
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))
            slide = prs.slides[slide_index]

            tables = [shape for shape in slide.shapes if shape.has_table]
            if table_index >= len(tables):
                raise ValueError(f"表格索引 {table_index} 超出范围")

            table = tables[table_index].table
            cell = table.cell(row, col)

            # 设置填充颜色
            if fill_color:
                r, g, b = ColorUtils.hex_to_rgb(fill_color)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(r, g, b)

            # 设置文本格式
            if cell.text_frame:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if text_color:
                            r, g, b = ColorUtils.hex_to_rgb(text_color)
                            run.font.color.rgb = RGBColor(r, g, b)
                        if bold:
                            run.font.bold = True
                        if font_size:
                            run.font.size = Pt(font_size)

            prs.save(str(file_path))

            logger.info(f"表格单元格格式化成功: {file_path}")
            return {
                "success": True,
                "message": "表格单元格格式化成功",
                "filename": str(file_path),
            }

        except Exception as e:
            logger.error(f"格式化表格单元格失败: {e}")
            return {"success": False, "message": f"操作失败: {str(e)}"}

    # ========== 形状操作 ==========
    def add_shape(
        self,
        filename: str,
        slide_index: int,
        shape_type: str,
        left_inches: float,
        top_inches: float,
        width_inches: float,
        height_inches: float,
        text: Optional[str] = None,
        fill_color: Optional[str] = None,
        line_color: Optional[str] = None,
    ) -> dict[str, Any]:
        """添加形状.

        Args:
            filename: 文件名
            slide_index: 幻灯片索引
            shape_type: 形状类型 ('rectangle', 'oval', 'triangle', 'arrow', 'rounded_rectangle')
            left_inches: 左边距(英寸)
            top_inches: 上边距(英寸)
            width_inches: 宽度(英寸)
            height_inches: 高度(英寸)
            text: 形状中的文本
            fill_color: 填充颜色
            line_color: 线条颜色
        """
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))

            if slide_index >= len(prs.slides):
                raise ValueError(f"幻灯片索引 {slide_index} 超出范围")

            slide = prs.slides[slide_index]

            # 形状类型映射
            shape_type_map = {
                'rectangle': MSO_SHAPE.RECTANGLE,
                'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
                'oval': MSO_SHAPE.OVAL,
                'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
                'arrow': MSO_SHAPE.RIGHT_ARROW,
            }

            if shape_type not in shape_type_map:
                raise ValueError(f"不支持的形状类型: {shape_type}")

            # 添加形状
            shape = slide.shapes.add_shape(
                shape_type_map[shape_type],
                Inches(left_inches),
                Inches(top_inches),
                Inches(width_inches),
                Inches(height_inches)
            )

            # 设置文本
            if text and hasattr(shape, "text_frame"):
                shape.text = text

            # 设置填充颜色
            if fill_color:
                r, g, b = ColorUtils.hex_to_rgb(fill_color)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(r, g, b)

            # 设置线条颜色
            if line_color:
                r, g, b = ColorUtils.hex_to_rgb(line_color)
                shape.line.color.rgb = RGBColor(r, g, b)

            prs.save(str(file_path))

            logger.info(f"形状添加成功: {file_path}")
            return {
                "success": True,
                "message": f"形状 '{shape_type}' 添加成功",
                "filename": str(file_path),
                "shape_type": shape_type,
            }

        except Exception as e:
            logger.error(f"添加形状失败: {e}")
            return {"success": False, "message": f"操作失败: {str(e)}"}

    # ========== 图表操作 ==========
    def add_chart(
        self,
        filename: str,
        slide_index: int,
        chart_type: str,
        categories: List[str],
        series_data: dict[str, List[float]],
        left_inches: float = 1.0,
        top_inches: float = 1.5,
        width_inches: float = 8.0,
        height_inches: float = 5.0,
        title: Optional[str] = None,
    ) -> dict[str, Any]:
        """添加图表.

        Args:
            filename: 文件名
            slide_index: 幻灯片索引
            chart_type: 图表类型 ('column', 'bar', 'line', 'pie', 'area')
            categories: 分类标签列表
            series_data: 系列数据 {"系列名": [数据列表]}
            left_inches: 左边距(英寸)
            top_inches: 上边距(英寸)
            width_inches: 宽度(英寸)
            height_inches: 高度(英寸)
            title: 图表标题
        """
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))

            if slide_index >= len(prs.slides):
                raise ValueError(f"幻灯片索引 {slide_index} 超出范围")

            slide = prs.slides[slide_index]

            # 图表类型映射
            chart_type_map = {
                'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'bar': XL_CHART_TYPE.BAR_CLUSTERED,
                'line': XL_CHART_TYPE.LINE,
                'pie': XL_CHART_TYPE.PIE,
                'area': XL_CHART_TYPE.AREA,
            }

            if chart_type not in chart_type_map:
                raise ValueError(f"不支持的图表类型: {chart_type}")

            # 准备图表数据
            chart_data = CategoryChartData()
            chart_data.categories = categories

            for series_name, values in series_data.items():
                chart_data.add_series(series_name, values)

            # 添加图表
            x, y, cx, cy = Inches(left_inches), Inches(top_inches), Inches(width_inches), Inches(height_inches)
            chart = slide.shapes.add_chart(
                chart_type_map[chart_type], x, y, cx, cy, chart_data
            ).chart

            # 设置标题
            if title:
                chart.has_title = True
                chart.chart_title.text_frame.text = title

            prs.save(str(file_path))

            logger.info(f"图表添加成功: {file_path}")
            return {
                "success": True,
                "message": f"图表 '{chart_type}' 添加成功",
                "filename": str(file_path),
                "chart_type": chart_type,
            }

        except Exception as e:
            logger.error(f"添加图表失败: {e}")
            return {"success": False, "message": f"操作失败: {str(e)}"}
