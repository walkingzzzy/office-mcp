"""PowerPoint 内容提取模块."""

from pathlib import Path
from typing import Any, Optional, List, Dict

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from loguru import logger

from office_mcp_server.config import config
from office_mcp_server.utils.file_manager import FileManager


class PowerPointContentExtraction:
    """PowerPoint 内容提取类."""

    def __init__(self) -> None:
        """初始化内容提取类."""
        self.file_manager = FileManager()

    def extract_all_text(self, filename: str) -> dict[str, Any]:
        """提取演示文稿中所有文本内容.
        
        Args:
            filename: 文件名
            
        Returns:
            dict: 包含所有文本内容的结果
        """
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)
            
            prs = Presentation(str(file_path))
            
            all_text = []
            slide_texts = []
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_content = {
                    "slide_index": slide_idx,
                    "texts": []
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_content["texts"].append(shape.text)
                        all_text.append(shape.text)
                
                slide_texts.append(slide_content)
            
            logger.info(f"文本提取成功: {file_path}, 共提取 {len(all_text)} 个文本块")
            return {
                "success": True,
                "message": "文本提取成功",
                "filename": str(file_path),
                "total_slides": len(prs.slides),
                "total_text_blocks": len(all_text),
                "slide_texts": slide_texts,
                "all_text": all_text
            }
            
        except Exception as e:
            logger.error(f"提取文本失败: {e}")
            return {"success": False, "message": f"提取失败: {str(e)}"}

    def extract_titles(self, filename: str) -> dict[str, Any]:
        """提取所有幻灯片标题.
        
        Args:
            filename: 文件名
            
        Returns:
            dict: 包含所有标题的结果
        """
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)
            
            prs = Presentation(str(file_path))
            
            titles = []
            
            for slide_idx, slide in enumerate(prs.slides):
                title_text = ""
                if slide.shapes.title:
                    title_text = slide.shapes.title.text
                
                titles.append({
                    "slide_index": slide_idx,
                    "title": title_text
                })
            
            logger.info(f"标题提取成功: {file_path}, 共提取 {len(titles)} 个标题")
            return {
                "success": True,
                "message": "标题提取成功",
                "filename": str(file_path),
                "total_slides": len(prs.slides),
                "titles": titles
            }
            
        except Exception as e:
            logger.error(f"提取标题失败: {e}")
            return {"success": False, "message": f"提取失败: {str(e)}"}

    def extract_notes(self, filename: str) -> dict[str, Any]:
        """提取所有演讲者备注.
        
        Args:
            filename: 文件名
            
        Returns:
            dict: 包含所有备注的结果
        """
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)
            
            prs = Presentation(str(file_path))
            
            notes = []
            
            for slide_idx, slide in enumerate(prs.slides):
                notes_text = ""
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    notes_text = notes_slide.notes_text_frame.text
                
                notes.append({
                    "slide_index": slide_idx,
                    "notes": notes_text
                })
            
            logger.info(f"备注提取成功: {file_path}, 共提取 {len(notes)} 个备注")
            return {
                "success": True,
                "message": "备注提取成功",
                "filename": str(file_path),
                "total_slides": len(prs.slides),
                "notes": notes
            }
            
        except Exception as e:
            logger.error(f"提取备注失败: {e}")
            return {"success": False, "message": f"提取失败: {str(e)}"}

    def extract_images(self, filename: str) -> dict[str, Any]:
        """提取图片信息列表.

        Args:
            filename: 文件名

        Returns:
            dict: 包含所有图片信息的结果
        """
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))

            images = []

            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_info = {
                            "slide_index": slide_idx,
                            "shape_index": shape_idx,
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height,
                            "name": shape.name if hasattr(shape, "name") else f"Picture {shape_idx}"
                        }

                        # 尝试获取图片文件信息
                        if hasattr(shape, "image"):
                            image_info["content_type"] = shape.image.content_type
                            image_info["ext"] = shape.image.ext

                        images.append(image_info)

            logger.info(f"图片信息提取成功: {file_path}, 共提取 {len(images)} 张图片")
            return {
                "success": True,
                "message": "图片信息提取成功",
                "filename": str(file_path),
                "total_images": len(images),
                "images": images
            }

        except Exception as e:
            logger.error(f"提取图片信息失败: {e}")
            return {"success": False, "message": f"提取失败: {str(e)}"}

    def extract_hyperlinks(self, filename: str) -> dict[str, Any]:
        """提取超链接列表.

        Args:
            filename: 文件名

        Returns:
            dict: 包含所有超链接的结果
        """
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))

            hyperlinks = []

            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    # 检查形状是否有超链接
                    if hasattr(shape, "click_action") and shape.click_action.hyperlink:
                        link_info = {
                            "slide_index": slide_idx,
                            "shape_index": shape_idx,
                            "shape_name": shape.name if hasattr(shape, "name") else f"Shape {shape_idx}",
                            "url": shape.click_action.hyperlink.address if shape.click_action.hyperlink.address else ""
                        }

                        # 获取链接文本
                        if hasattr(shape, "text"):
                            link_info["text"] = shape.text

                        hyperlinks.append(link_info)

                    # 检查文本框中的超链接
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if hasattr(run, "hyperlink") and run.hyperlink and run.hyperlink.address:
                                    link_info = {
                                        "slide_index": slide_idx,
                                        "shape_index": shape_idx,
                                        "shape_name": shape.name if hasattr(shape, "name") else f"Shape {shape_idx}",
                                        "text": run.text,
                                        "url": run.hyperlink.address
                                    }
                                    hyperlinks.append(link_info)

            logger.info(f"超链接提取成功: {file_path}, 共提取 {len(hyperlinks)} 个超链接")
            return {
                "success": True,
                "message": "超链接提取成功",
                "filename": str(file_path),
                "total_links": len(hyperlinks),
                "hyperlinks": hyperlinks
            }

        except Exception as e:
            logger.error(f"提取超链接失败: {e}")
            return {"success": False, "message": f"提取失败: {str(e)}"}

    def extract_all_content(self, filename: str) -> dict[str, Any]:
        """提取演示文稿的所有内容（文本、标题、备注、图片、超链接）.

        Args:
            filename: 文件名

        Returns:
            dict: 包含所有内容的综合结果
        """
        try:
            # 调用各个提取方法
            text_result = self.extract_all_text(filename)
            titles_result = self.extract_titles(filename)
            notes_result = self.extract_notes(filename)
            images_result = self.extract_images(filename)
            links_result = self.extract_hyperlinks(filename)

            # 检查是否所有提取都成功
            if not all([
                text_result.get("success"),
                titles_result.get("success"),
                notes_result.get("success"),
                images_result.get("success"),
                links_result.get("success")
            ]):
                return {
                    "success": False,
                    "message": "部分内容提取失败"
                }

            logger.info(f"完整内容提取成功: {filename}")
            return {
                "success": True,
                "message": "完整内容提取成功",
                "filename": filename,
                "text": text_result,
                "titles": titles_result,
                "notes": notes_result,
                "images": images_result,
                "hyperlinks": links_result
            }

        except Exception as e:
            logger.error(f"提取完整内容失败: {e}")
            return {"success": False, "message": f"提取失败: {str(e)}"}

