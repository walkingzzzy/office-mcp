"""PowerPoint 导出操作模块."""

from typing import Any, Optional

from pptx import Presentation
from loguru import logger

from office_mcp_server.config import config
from office_mcp_server.utils.file_manager import FileManager


class PowerPointExportOperations:
    """PowerPoint 导出操作类."""

    def __init__(self) -> None:
        """初始化导出操作类."""
        self.file_manager = FileManager()

    def export_presentation(
        self,
        filename: str,
        export_format: str = "pdf",
        output_filename: Optional[str] = None,
    ) -> dict[str, Any]:
        """导出PowerPoint演示文稿到其他格式."""
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))

            # 确定输出文件名
            if not output_filename:
                base_name = file_path.stem
                if export_format == 'pdf':
                    output_filename = f"{base_name}.pdf"
                elif export_format == 'html':
                    output_filename = f"{base_name}.html"
                elif export_format == 'images':
                    output_filename = f"{base_name}_slides"
                else:
                    raise ValueError(f"不支持的导出格式: {export_format}")

            output_path = config.paths.output_dir / output_filename

            # 根据格式导出
            if export_format == 'pdf':
                # PDF导出
                try:
                    import comtypes.client

                    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                    powerpoint.Visible = 1

                    deck = powerpoint.Presentations.Open(str(file_path.absolute()))
                    deck.SaveAs(str(output_path.absolute()), 32)  # 32 = ppSaveAsPDF
                    deck.Close()
                    powerpoint.Quit()

                    message = f"演示文稿已成功导出为 PDF: {output_path}"
                except ImportError:
                    return {
                        "success": False,
                        "message": "PDF导出需要 Microsoft PowerPoint 或安装 comtypes 库"
                    }
                except Exception as e:
                    return {
                        "success": False,
                        "message": f"PDF导出失败（请确保已安装 Microsoft PowerPoint）: {str(e)}"
                    }

            elif export_format == 'html':
                # 导出为HTML
                html_content = self._convert_to_html(prs)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(html_content)
                message = f"演示文稿已成功导出为 HTML: {output_path}"

            elif export_format == 'images':
                # 导出为图片序列
                try:
                    from PIL import Image
                    import io

                    output_dir = config.paths.output_dir / output_filename
                    output_dir.mkdir(exist_ok=True)

                    message = f"图片导出功能需要额外的图像处理库。已创建目录: {output_dir}"

                    return {
                        "success": False,
                        "message": "图片导出功能需要系统级PowerPoint支持或专门的转换工具"
                    }
                except ImportError:
                    return {
                        "success": False,
                        "message": "图片导出需要安装 Pillow 库: pip install Pillow"
                    }

            else:
                raise ValueError(f"不支持的导出格式: {export_format}")

            logger.info(f"演示文稿导出成功: {output_path}")
            return {
                "success": True,
                "message": message,
                "source_file": str(file_path),
                "output_file": str(output_path),
                "format": export_format,
            }

        except Exception as e:
            logger.error(f"导出演示文稿失败: {e}")
            return {"success": False, "message": f"导出失败: {str(e)}"}

    def _convert_to_html(self, prs: Presentation) -> str:
        """将PowerPoint演示文稿转换为HTML."""
        html_parts = ['<!DOCTYPE html>', '<html>', '<head>',
                     '<meta charset="UTF-8">',
                     '<title>演示文稿</title>',
                     '<style>',
                     '.slide { page-break-after: always; margin: 20px; border: 1px solid #ccc; padding: 20px; }',
                     '.slide-title { font-size: 24px; font-weight: bold; margin-bottom: 10px; }',
                     '.slide-content { font-size: 16px; }',
                     '</style>',
                     '</head>', '<body>']

        for i, slide in enumerate(prs.slides):
            html_parts.append(f'<div class="slide" id="slide-{i+1}">')
            html_parts.append(f'<h2>幻灯片 {i+1}</h2>')

            # 提取文本内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if hasattr(shape, "text_frame"):
                        if shape.text_frame.text and len(shape.text_frame.text) < 100:
                            html_parts.append(f'<div class="slide-title">{shape.text}</div>')
                        else:
                            html_parts.append(f'<div class="slide-content">{shape.text}</div>')

            html_parts.append('</div>')

        html_parts.extend(['</body>', '</html>'])
        return '\n'.join(html_parts)
