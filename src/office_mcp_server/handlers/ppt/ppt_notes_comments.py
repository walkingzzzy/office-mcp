"""PowerPoint 备注和批注模块."""

from typing import Any, Optional

from pptx import Presentation
from loguru import logger

from office_mcp_server.config import config
from office_mcp_server.utils.file_manager import FileManager


class PowerPointNotesCommentsOperations:
    """PowerPoint 备注和批注操作类."""

    def __init__(self) -> None:
        """初始化备注和批注操作类."""
        self.file_manager = FileManager()

    def add_speaker_notes(
        self,
        filename: str,
        slide_index: int,
        notes_text: str,
    ) -> dict[str, Any]:
        """添加演讲者备注.

        Args:
            filename: 文件名
            slide_index: 幻灯片索引
            notes_text: 备注文本
        """
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))

            if slide_index >= len(prs.slides):
                raise ValueError(f"幻灯片索引 {slide_index} 超出范围")

            slide = prs.slides[slide_index]

            # 获取或创建备注页
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes_text

            prs.save(str(file_path))

            logger.info(f"演讲者备注添加成功: {file_path}")
            return {
                "success": True,
                "message": "演讲者备注添加成功",
                "filename": str(file_path),
                "slide_index": slide_index,
            }

        except Exception as e:
            logger.error(f"添加演讲者备注失败: {e}")
            return {"success": False, "message": f"操作失败: {str(e)}"}

    def get_speaker_notes(
        self,
        filename: str,
        slide_index: int,
    ) -> dict[str, Any]:
        """获取演讲者备注.

        Args:
            filename: 文件名
            slide_index: 幻灯片索引
        """
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            prs = Presentation(str(file_path))

            if slide_index >= len(prs.slides):
                raise ValueError(f"幻灯片索引 {slide_index} 超出范围")

            slide = prs.slides[slide_index]
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text

            logger.info(f"演讲者备注获取成功: {file_path}")
            return {
                "success": True,
                "message": "演讲者备注获取成功",
                "filename": str(file_path),
                "slide_index": slide_index,
                "notes_text": notes_text,
            }

        except Exception as e:
            logger.error(f"获取演讲者备注失败: {e}")
            return {"success": False, "message": f"操作失败: {str(e)}"}

    def add_comment(
        self,
        filename: str,
        slide_index: int,
        author: str,
        comment_text: str,
        x_pos: float = 0.0,
        y_pos: float = 0.0,
    ) -> dict[str, Any]:
        """添加批注.

        Args:
            filename: 文件名
            slide_index: 幻灯片索引
            author: 作者名称
            comment_text: 批注文本
            x_pos: X坐标位置
            y_pos: Y坐标位置

        Note:
            python-pptx库对批注的支持有限,完整功能需要使用win32com
        """
        try:
            file_path = config.paths.output_dir / filename
            self.file_manager.validate_file_path(file_path, must_exist=True)

            # python-pptx不直接支持添加批注
            # 需要使用win32com或更底层的XML操作
            return {
                "success": False,
                "message": "python-pptx库不支持添加批注,需要使用win32com或更底层的XML操作",
                "filename": str(file_path),
                "alternative": "可以使用演讲者备注作为替代方案",
            }

        except Exception as e:
            logger.error(f"添加批注失败: {e}")
            return {"success": False, "message": f"操作失败: {str(e)}"}
